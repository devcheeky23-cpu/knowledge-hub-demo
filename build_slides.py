from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Theme
# ---------------------------------------------------------------------------
BG      = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
ACCENT  = RGBColor(0x38, 0xBD, 0xF8)   # sky blue
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0x94, 0xA3, 0xB8)
GREEN   = RGBColor(0x34, 0xD3, 0x99)
YELLOW  = RGBColor(0xFB, 0xBF, 0x24)
RED     = RGBColor(0xF8, 0x71, 0x71)

W = Inches(13.33)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def bg(slide, color=BG):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def txt(slide, text, l, t, w, h, size=24, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def accent_bar(slide, t=Inches(1.1), h=Inches(0.05)):
    bar = slide.shapes.add_shape(1, Inches(0.6), t, Inches(12.1), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def step_tag(slide, n, label):
    """Small '1 · PROBLEM' style tag above the title."""
    txt(slide, f"{n}  ·  {label}", Inches(0.62), Inches(0.32), Inches(8), Inches(0.4),
        size=14, bold=True, color=ACCENT)


def bullet_box(slide, items, l, t, w, h, size=20, color=WHITE, marker="•"):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"{marker}  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color


def colored_box(slide, label, desc, l, t, w, h, box_color):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = box_color
    shape.line.fill.background()
    txt(slide, label, l + Inches(0.15), t + Inches(0.1), w - Inches(0.3),
        Inches(0.4), size=16, bold=True, color=BG)
    txt(slide, desc,  l + Inches(0.15), t + Inches(0.45), w - Inches(0.3),
        h - Inches(0.55), size=13, color=BG, wrap=True)


def table_2col(slide, header_l, header_r, rows, l, t, w, size=15,
               accent=ACCENT):
    """Simple two-column table rendered with textboxes."""
    col_w = w // 2
    # headers
    txt(slide, header_l, l, t, col_w, Inches(0.4), size=size + 1, bold=True, color=accent)
    txt(slide, header_r, l + col_w, t, col_w, Inches(0.4), size=size + 1, bold=True, color=accent)
    row_h = Inches(0.62)
    for i, (a, b) in enumerate(rows):
        y = t + Inches(0.5) + i * row_h
        txt(slide, a, l, y, col_w - Inches(0.1), row_h, size=size, color=WHITE)
        txt(slide, b, l + col_w, y, col_w - Inches(0.1), row_h, size=size, color=GRAY)


# ===========================================================================
# Slides — Executive Storyline structure
#   Title → Problem → Value → Approach+Demo → Risk & Control →
#   Implementation Path → Q&A Buffer (appendix)
# ===========================================================================

prs = new_prs()

# ---------------------------------------------------------------------------
# 1. Title
# ---------------------------------------------------------------------------
s = blank_slide(prs)
bg(s)
txt(s, "Project Knowledge Hub", Inches(0.6), Inches(1.8), Inches(12), Inches(1.2),
    size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Ask your project documents in natural language — get cited, trustworthy answers",
    Inches(0.6), Inches(3.1), Inches(12), Inches(0.9),
    size=22, color=ACCENT, align=PP_ALIGN.CENTER)
txt(s, "Advanced AI Engineering  ·  BOI STEM++",
    Inches(0.6), Inches(4.3), Inches(12), Inches(0.5),
    size=18, color=GRAY, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------------------
# 2. Problem
# ---------------------------------------------------------------------------
s = blank_slide(prs)
bg(s)
step_tag(s, "1", "PROBLEM")
txt(s, "Project knowledge is scattered — and it costs us", Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
    size=32, bold=True, color=WHITE)
accent_bar(s, t=Inches(1.45))

bullet_box(s, [
    "Developers repeatedly ask the same questions — API fields, design decisions, requirements",
    "Knowledge hand-off falls through the cracks: PM / BA / SA → dev, and backend → frontend",
    "When someone leaves, their knowledge leaves with them",
    "No single source of truth — it lives in docs, chat threads, and people's heads",
], Inches(0.6), Inches(1.8), Inches(12), Inches(3.2), size=21)

txt(s, "Impact:  slower delivery  ·  long onboarding  ·  repeated Q&A drains the team every week",
    Inches(0.6), Inches(5.6), Inches(12), Inches(0.8),
    size=19, color=YELLOW)

# ---------------------------------------------------------------------------
# 3. Value
# ---------------------------------------------------------------------------
s = blank_slide(prs)
bg(s)
step_tag(s, "2", "VALUE")
txt(s, "One place to ask — answers from real documents", Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
    size=32, bold=True, color=WHITE)
accent_bar(s, t=Inches(1.45))

table_2col(s, "Outcome", "How we measure it", [
    ("Devs self-serve answers, stop interrupting others", "Fewer repeated questions / hours saved"),
    ("New hires get productive faster", "Time-to-productive for onboarding"),
    ("Knowledge survives staff turnover", "Questions answerable from documents"),
    ("Surface where documentation is missing", "Gap report (questions the system abstains on)"),
], Inches(0.6), Inches(1.9), Inches(12.1), size=18)

# ---------------------------------------------------------------------------
# 4a. Approach
# ---------------------------------------------------------------------------
s = blank_slide(prs)
bg(s)
step_tag(s, "3", "APPROACH")
txt(s, "RAG: answer only from our documents, always cite", Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
    size=30, bold=True, color=WHITE)
accent_bar(s, t=Inches(1.45))

# Query flow diagram
for i, (label, x) in enumerate([
    ("User\nQuestion", 0.6),
    ("Retrieve\nrelevant chunks", 3.0),
    ("Augment prompt\n+ instructions", 5.4),
    ("LLM\nGenerate", 7.8),
    ("Answer +\nCitation", 10.2),
]):
    box = s.shapes.add_shape(1, Inches(x), Inches(1.9), Inches(2.1), Inches(1.0))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x5F)
    box.line.color.rgb = ACCENT
    txt(s, label, Inches(x + 0.1), Inches(1.95), Inches(1.9), Inches(0.9),
        size=13, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        txt(s, "→", Inches(x + 2.15), Inches(2.2), Inches(0.5), Inches(0.5),
            size=20, color=ACCENT, align=PP_ALIGN.CENTER)

txt(s, "Data & boundary", Inches(0.6), Inches(3.4), Inches(12), Inches(0.5),
    size=18, bold=True, color=ACCENT)
bullet_box(s, [
    "Answers come only from documents imported into the system (one project)",
    "The model is constrained to the retrieved context — it does not use general knowledge",
    "Every answer carries a citation — the user can verify the source",
], Inches(0.6), Inches(3.9), Inches(12), Inches(2.2), size=19)

# ---------------------------------------------------------------------------
# 4b. Demo
# ---------------------------------------------------------------------------
s = blank_slide(prs)
bg(s)
step_tag(s, "3", "APPROACH  ·  LIVE DEMO")
txt(s, "Three behaviors that make it trustworthy", Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
    size=30, bold=True, color=WHITE)
accent_bar(s, t=Inches(1.45))

scenarios = [
    (GREEN,  "✅  Found",
     'Q: "What fields does the Order API return?"\n→  Answer + clickable citation\n     (api-spec.md › Order Endpoints)'),
    (YELLOW, "⚠️  Abstain",
     'Q: "Does the system support GDPR?"\n→  "ไม่พบข้อมูลในเอกสาร"\n     No guessing.'),
    (RED,    "⚡  Conflict",
     'Q: "What is the payment timeout?"\n→  api-spec.md says 30s\n     architecture.md says 60s\n     Both shown. No winner picked.'),
]
for i, (color, title, body) in enumerate(scenarios):
    x = Inches(0.5 + i * 4.25)
    box = s.shapes.add_shape(1, x, Inches(1.8), Inches(4.0), Inches(4.0))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x1A, 0x22, 0x3A)
    box.line.color.rgb = color
    txt(s, title, x + Inches(0.15), Inches(1.9), Inches(3.7), Inches(0.5),
        size=18, bold=True, color=color)
    txt(s, body, x + Inches(0.15), Inches(2.5), Inches(3.7), Inches(3.2),
        size=15, color=WHITE, wrap=True)

txt(s, 'Design principle:  "abstention is a feature, not a failure" — a system that guesses confidently is more dangerous than one that says "I don\'t know"',
    Inches(0.6), Inches(6.1), Inches(12), Inches(0.8), size=15, color=GRAY)

# ---------------------------------------------------------------------------
# 5. Risk & Control
# ---------------------------------------------------------------------------
s = blank_slide(prs)
bg(s)
step_tag(s, "4", "RISK & CONTROL")
txt(s, "Risk doesn't mean forbidden — it means controlled", Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
    size=30, bold=True, color=WHITE)
accent_bar(s, t=Inches(1.45))

# three-column header
headers = ["Risk", "Control", "Owner"]
xs = [0.6, 6.0, 11.0]
ws = [5.3, 4.9, 1.8]
for hx, hw, htext in zip(xs, ws, headers):
    txt(s, htext, Inches(hx), Inches(1.7), Inches(hw), Inches(0.4),
        size=16, bold=True, color=ACCENT)

rows = [
    ("Hallucination (made-up answers)", "Answer-from-docs only + forced citation + abstain", "Dev team"),
    ("Data leakage", "MVP uses mock data; production → paid/self-host", "Data owner"),
    ("Wrong info used in production", "Every answer cited — user verifies before use", "User"),
    ("Access boundary", "One project; future: per-user access control", "Admin"),
    ("Conflicting documents", "System reports conflict, does not decide", "User"),
]
for i, (a, b, c) in enumerate(rows):
    y = Inches(2.2) + i * Inches(0.78)
    txt(s, a, Inches(0.6), y, Inches(5.3), Inches(0.75), size=14, color=WHITE)
    txt(s, b, Inches(6.0), y, Inches(4.9), Inches(0.75), size=14, color=GRAY)
    txt(s, c, Inches(11.0), y, Inches(1.8), Inches(0.75), size=14, color=GREEN)

# ---------------------------------------------------------------------------
# 6. Implementation Path
# ---------------------------------------------------------------------------
s = blank_slide(prs)
bg(s)
step_tag(s, "5", "IMPLEMENTATION PATH")
txt(s, "Pilot in 2 weeks, then scale", Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
    size=32, bold=True, color=WHITE)
accent_bar(s, t=Inches(1.45))

# Pilot column
header = s.shapes.add_shape(1, Inches(0.5), Inches(1.8), Inches(6.0), Inches(0.5))
header.fill.solid(); header.fill.fore_color.rgb = ACCENT; header.line.fill.background()
txt(s, "Pilot  ·  2 weeks", Inches(0.6), Inches(1.85), Inches(5.8), Inches(0.45),
    size=17, bold=True, color=BG)
bullet_box(s, [
    "One project, seed 5–10 real documents (no sensitive data)",
    "Build a golden question set (15–20) from questions devs actually ask",
    "Measure: ≥70–80% correct + every answer cited + correct abstention",
], Inches(0.55), Inches(2.45), Inches(6.1), Inches(3.5), size=17)

# Scale column
header = s.shapes.add_shape(1, Inches(6.9), Inches(1.8), Inches(6.0), Inches(0.5))
header.fill.solid(); header.fill.fore_color.rgb = GREEN; header.line.fill.background()
txt(s, "Scale", Inches(7.0), Inches(1.85), Inches(5.8), Inches(0.45),
    size=17, bold=True, color=BG)
bullet_box(s, [
    "Move to persistent storage + paid / self-hosted model (privacy)",
    "Support multiple projects + per-user access control",
    "Add gap report, code indexing, auto-sync with git",
], Inches(6.95), Inches(2.45), Inches(6.0), Inches(3.5), size=17)

txt(s, "MVP runs at zero cost today. Production investment is a deliberate, measured next step.",
    Inches(0.6), Inches(6.2), Inches(12), Inches(0.6), size=16, color=ACCENT)

# ---------------------------------------------------------------------------
# 7. Q&A Buffer (appendix — not presented, used to answer questions)
# ---------------------------------------------------------------------------
s = blank_slide(prs)
bg(s)
step_tag(s, "+", "APPENDIX  ·  Q&A BUFFER")
txt(s, "Anticipated questions", Inches(0.6), Inches(0.7), Inches(12), Inches(0.7),
    size=30, bold=True, color=WHITE)
accent_bar(s, t=Inches(1.45))

qa = [
    ("How is this different from ChatGPT?",
     "ChatGPT can guess, sends data outside the org, and you can't control its behavior. Ours answers only from our documents, always cites, lets us control retrieval / abstain / conflict, and can be deployed in our own infrastructure."),
    ("Can't ChatGPT also abstain if instructed?",
     "Yes — but for us it's a guaranteed default, it's measurable, and we control retrieval ourselves (including Thai-language matching). ChatGPT's file upload doesn't expose those knobs."),
    ("Then why not just use ChatGPT?",
     "Privacy + retrieval control (Thai) + conflict mode + the ability to integrate into our own workflow and infrastructure."),
]
y = 1.9
for q, a in qa:
    txt(s, f"Q  {q}", Inches(0.6), Inches(y), Inches(12), Inches(0.5),
        size=18, bold=True, color=ACCENT)
    txt(s, f"A  {a}", Inches(0.6), Inches(y + 0.5), Inches(12), Inches(1.1),
        size=15, color=WHITE)
    y += 1.75


prs.save("knowledge-hub-pitch.pptx")
print("Saved: knowledge-hub-pitch.pptx  (7 slides, Executive Storyline structure)")
